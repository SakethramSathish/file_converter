import streamlit as st
from pptx import Presentation
import os
import tempfile
import comtypes.client

def convert():
    option = st.selectbox("Choose a Conversion", ["PPTX to PDF", "PPTX to Images", "PPTX to Text"])

    file = st.file_uploader("Upload a PowerPoint File", type=["pptx"])

    if file and st.button("Convert"):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as input_tmp:
            input_tmp.write(file.read())
            input_path = input_tmp.name

        base_name = os.path.splitext(input_path)[0]
        try:
            if option == "PPTX to PDF":
                output_file = input_path.replace(".pptx", ".pdf")
                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                powerpoint.Visible = 1
                presentation = powerpoint.Presentations.Open(input_path)
                presentation.SaveAs(output_file, 32)
                presentation.Close()
                powerpoint.Quit()

                with open(output_file, "rb") as f:
                    st.download_button("Download PDF File", f, file_name = f"{base_name}.pdf")

            elif option == "PPTX to Images":
                prs = Presentation(input_path)
                text = ""
                for i, slide in enumerate(prs.slides):
                    img_text = f"Slide {i + 1}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            img_text += shape.text + "\n"
                    st.text_area(f"Text from Slide {i + 1}", img_text.strip(), height = 150)
            
            elif option == "PPTX to Text":
                prs = Presentation(input_path)
                all_text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            all_text += shape.text + "\n"
                st.download_button("Download Text File", all_text, file_name = f"{base_name}.txt")

            else:
                st.error("Unsupported conversion type.")

        except Exception as e:
            st.error(f"Conversion failed: {str(e)}")
                